"""
Slide-deck ingestion worker.

Extracts each slide's body text and speaker notes from a .pptx, writing the
combined text to the registry (so it indexes into media_text) and the per-slide
structure to asset.meta. Thumbnails need a renderer (LibreOffice headless) and
are deferred; text is the searchable payload for now.

python-pptx is imported lazily and a missing one degrades to a registered-but-
empty asset rather than a failure. .key (Keynote) is not parsed here.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from .. import registry


def _extract_pptx(path: Path) -> list[dict]:
    """Return [{no, text, notes}, ...] per slide, or [] if python-pptx is absent."""
    try:
        from pptx import Presentation
    except Exception:
        return []
    slides: list[dict] = []
    try:
        prs = Presentation(str(path))
    except Exception:
        return []
    for i, slide in enumerate(prs.slides, start=1):
        body = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    body.append(t)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes = ""
        slides.append({"no": i, "text": "\n".join(body), "notes": notes})
    return slides


def enrich(asset_id: str, path: str) -> dict[str, Any]:
    deck_path = Path(path)
    if not deck_path.exists():
        registry.set_status(asset_id, "failed")
        return {"status": "failed", "detail": "file not found"}

    slides = _extract_pptx(deck_path) if deck_path.suffix.lower() == ".pptx" else []
    notes_list: list[str] = []
    if not slides:
        notes_list.append("no slide text (python-pptx missing, .key, or empty)")

    # Combined text, slide markers preserved so a chunk reads in context.
    parts = []
    for s in slides:
        block = f"Slide {s['no']}: {s['text']}".strip()
        if s["notes"]:
            block += f"\nNotes: {s['notes']}"
        parts.append(block)
        registry.add_moment(
            asset_id,
            kind="slide",
            label=f"Slide {s['no']}",
            text=block,
            meta={"slide_no": s["no"], "notes": s["notes"]},
        )
    combined = "\n\n".join(parts)

    if combined.strip():
        registry.add_transcript(asset_id, language=None, segments=[], text=combined)

    registry.update_asset(
        asset_id,
        status="ready",
        meta={"slide_count": len(slides),
              "slides": slides,
              "thumbnails": "deferred (needs LibreOffice render)"},
    )
    return {
        "status":      "ready",
        "slide_count": len(slides),
        "chars":       len(combined),
        "notes":       notes_list,
    }
